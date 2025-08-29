import io
from typing import Dict, List, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt

def collect_template_images(prs: Presentation) -> List[bytes]:
    """Collect raw image blobs from picture shapes in the template deck to optionally reuse."""
    images = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE (avoid import for simplicity)
                try:
                    images.append(shape.image.blob)
                except Exception:
                    pass
    return images

def add_bulleted_slide(prs: Presentation, layout_idx: int, title: str, bullets: List[str], notes: str, image_blobs: List[bytes], image_cycle_idx: int) -> int:
    """Add a slide using the given layout index, populate title and bullets, optionally reuse an image."""
    slide_layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(slide_layout)
    # Try to place title in first title placeholder (if any)
    try:
        slide.shapes.title.text = title or ""
    except Exception:
        pass

    # Try to find a content placeholder for bullets
    content_placeholder = None
    for shp in slide.placeholders:
        # placeholder type 1 = TITLE, 2 = BODY
        if getattr(shp, "placeholder_format", None):
            if shp.placeholder_format.type == 2:
                content_placeholder = shp
                break
    if content_placeholder:
        tf = content_placeholder.text_frame
        tf.clear()
        if bullets:
            tf.text = bullets[0]
            for b in bullets[1:]:
                p = tf.add_paragraph()
                p.text = b
                p.level = 0

    # Optionally drop in an image from the template if the layout has free space.
    # We'll place it at bottom-right with a small size, harmless if overlapping.
    if image_blobs:
        blob = image_blobs[image_cycle_idx % len(image_blobs)]
        try:
            slide.shapes.add_picture(io.BytesIO(blob), Inches(8.0), Inches(4.5), width=Inches(2.0))
        except Exception:
            pass

    # Speaker notes
    if notes:
        try:
            slide.notes_slide.notes_text_frame.text = notes
        except Exception:
            pass

    return image_cycle_idx + 1

def build_presentation(template_bytes: bytes, deck: Dict, prefer_layout_idx: int = 1) -> bytes:
    """
    Create a new .pptx from the uploaded template and a structured deck dict:
    {"title": str, "slides": [{"title":..., "bullets":[...], "notes": "..."}]}
    """
    # Start from template to inherit theme, colors, fonts, and layouts
    in_buf = io.BytesIO(template_bytes)
    prs = Presentation(in_buf)
    template_images = collect_template_images(prs)

    # If deck has an overall title, try to create a title slide (layout 0 is usually title)
    if deck.get("title"):
        try:
            title_slide_layout = prs.slide_layouts[0]
            ts = prs.slides.add_slide(title_slide_layout)
            try:
                ts.shapes.title.text = deck["title"]
            except Exception:
                pass
        except Exception:
            pass

    image_idx = 0
    slides = deck.get("slides", [])
    # Choose a reasonable default layout index for bullet slides: commonly 1 or 1-based "Title and Content"
    default_layout_idx = prefer_layout_idx if prefer_layout_idx < len(prs.slide_layouts) else 1

    for s in slides:
        title = s.get("title", "")
        bullets = [b for b in s.get("bullets", []) if b]
        notes = s.get("notes", "")
        image_idx = add_bulleted_slide(prs, default_layout_idx, title, bullets, notes, template_images, image_idx)

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()
